from __future__ import annotations

import base64
from pathlib import Path

import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient
from pptx import Presentation

from backend.api.routes.slide_decks import get_slide_deck_service, router as slide_decks_router
from backend.core.services.slide_deck_service import SlideDeckService
from backend.dependencies import get_knowledge_repository
from backend.domain.slide_deck import (
    SlideAsset,
    SlideAssetStage,
    SlideDeckFileKind,
    SlideDeckProject,
    SlideDeckStage,
    SlideDeckStatus,
    SlideRecord,
    SlideStatus,
)
from backend.infrastructure.repositories.memory_repository import InMemoryKnowledgeRepository
from backend.infrastructure.slide_deck_files import SlideDeckFileStore


ONE_PIXEL_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAFgwJ/lwWqJwAAAABJRU5ErkJggg=="
)


class FakeImageProvider:
    async def generate_image(self, prompt: str, aspect_ratio: str = "16:9", quality: str = "2K"):
        return type(
            "ImageResult",
            (),
            {"base64_data": base64.b64encode(ONE_PIXEL_PNG).decode(), "mime_type": "image/png"},
        )()


def build_export_client(tmp_path: Path):
    repo = InMemoryKnowledgeRepository()
    file_store = SlideDeckFileStore(tmp_path / "output")
    deck = SlideDeckProject(
        title="Exportable Deck",
        source_ids=["src_alpha"],
        status=SlideDeckStatus.READY,
        stage=SlideDeckStage.SLIDES_READY,
        config_snapshot={"aspect_ratio": "16:9"},
    )

    async def seed():
        for page in range(1, 3):
            stored = file_store.save_file(
                deck_id=deck.id,
                kind=SlideDeckFileKind.SLIDE_IMAGE,
                filename=f"slide_{page}.png",
                content=ONE_PIXEL_PNG,
            )
            slide = SlideRecord(
                deck_id=deck.id,
                page_number=page,
                title=f"Slide {page}",
                status=SlideStatus.SUCCEEDED,
            )
            asset = SlideAsset(
                deck_id=deck.id,
                slide_id=slide.id,
                file_path=str(stored.path),
                mime_type="image/png",
                byte_size=stored.byte_size,
                checksum=stored.checksum,
                download_ref=stored.download_ref,
                stage=SlideAssetStage.GENERATED,
            )
            await repo.save_slide_asset(asset)
            slide.asset_id = asset.id
            deck.slides.append(slide)
        await repo.save_slide_deck(deck)

    import asyncio

    asyncio.run(seed())

    service = SlideDeckService(
        repository=repo,
        planning_service=None,
        image_provider=FakeImageProvider(),
        edit_provider=None,
        file_store=file_store,
    )
    app = FastAPI()
    app.include_router(slide_decks_router, prefix="/api")
    app.dependency_overrides[get_slide_deck_service] = lambda: service
    app.dependency_overrides[get_knowledge_repository] = lambda: repo
    return TestClient(app), repo, deck


def test_slide_deck_export_job_creates_downloadable_image_based_pptx(tmp_path: Path):
    client, repo, deck = build_export_client(tmp_path)

    job = client.post(f"/api/slide-decks/{deck.id}/export/jobs")

    assert job.status_code == 200
    assert job.json()["status"] == "succeeded"
    export_id = job.json()["result_ref"]
    export = repo.slide_exports[export_id]
    assert export.format == "pptx"
    assert export.status == "succeeded"
    assert export.slide_count == 2
    assert export.filename.endswith(".pptx")
    assert Path(export.file_path).exists()

    presentation = Presentation(export.file_path)
    assert len(presentation.slides) == 2

    download = client.get(f"/api/slide-decks/{deck.id}/download?format=pptx")
    assert download.status_code == 200
    assert download.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert download.content == Path(export.file_path).read_bytes()

    image = client.get(f"/api/slide-decks/{deck.id}/slides/{deck.slides[0].id}/image")
    assert image.status_code == 200
    assert image.headers["content-type"].startswith("image/png")
    assert image.content == ONE_PIXEL_PNG


def test_slide_deck_export_preserves_four_by_three_dimensions(tmp_path: Path):
    client, repo, deck = build_export_client(tmp_path)
    deck.config_snapshot["aspect_ratio"] = "4:3"
    repo.slide_decks[deck.id] = deck

    job = client.post(f"/api/slide-decks/{deck.id}/export/jobs")

    assert job.status_code == 200
    export = repo.slide_exports[job.json()["result_ref"]]
    presentation = Presentation(export.file_path)
    assert presentation.slide_width == 9144000
    assert presentation.slide_height == 6858000


def test_slide_deck_export_rejects_deck_without_generated_slide_images(tmp_path: Path):
    client, repo, deck = build_export_client(tmp_path)
    deck.slides[0].asset_id = None
    repo.slide_decks[deck.id] = deck

    job = client.post(f"/api/slide-decks/{deck.id}/export/jobs")

    assert job.status_code == 200
    assert job.json()["status"] == "failed"
    assert "generated slide images" in job.json()["error"]


def test_slide_deck_download_requires_successful_export(tmp_path: Path):
    client, _, deck = build_export_client(tmp_path)

    download = client.get(f"/api/slide-decks/{deck.id}/download?format=pptx")

    assert download.status_code == 404
    assert "export" in download.json()["detail"]


def test_slide_deck_download_rejects_export_after_slide_regeneration(tmp_path: Path):
    client, repo, deck = build_export_client(tmp_path)
    export_job = client.post(f"/api/slide-decks/{deck.id}/export/jobs")
    assert export_job.status_code == 200
    assert export_job.json()["status"] == "succeeded"
    old_export = repo.slide_exports[export_job.json()["result_ref"]]
    old_asset_id = deck.slides[0].asset_id

    regenerated = client.post(f"/api/slide-decks/{deck.id}/slides/{deck.slides[0].id}/regenerate")
    assert regenerated.status_code == 200
    assert regenerated.json()["slides"][0]["asset_id"] != old_asset_id

    download = client.get(f"/api/slide-decks/{deck.id}/download?format=pptx")

    assert download.status_code == 404
    assert Path(old_export.file_path).exists()


def test_slide_deck_download_returns_clean_errors_for_missing_deck_and_unsupported_format(tmp_path: Path):
    client, _, deck = build_export_client(tmp_path)

    unsupported = client.get(f"/api/slide-decks/{deck.id}/download?format=pdf")
    missing = client.get("/api/slide-decks/deck_missing/download?format=pptx")

    assert unsupported.status_code == 400
    assert "Unsupported export format" in unsupported.json()["detail"]
    assert missing.status_code == 404
    assert "slide deck not found" in missing.json()["detail"]
