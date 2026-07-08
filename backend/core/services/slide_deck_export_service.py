"""PPTX export helpers for image-based slide decks."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path


class SlideDeckPPTXExporter:
    def export(self, image_paths: list[str], aspect_ratio: str = "16:9") -> bytes:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as exc:
            raise RuntimeError("python-pptx is required for PPTX export") from exc

        presentation = Presentation()
        if aspect_ratio == "4:3":
            presentation.slide_width = Inches(10)
            presentation.slide_height = Inches(7.5)
        else:
            presentation.slide_width = Inches(10)
            presentation.slide_height = Inches(5.625)

        blank_layout = presentation.slide_layouts[6]
        for image_path in image_paths:
            if not Path(image_path).exists():
                raise ValueError(f"slide image does not exist: {image_path}")
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                image_path,
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )

        buffer = BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()
